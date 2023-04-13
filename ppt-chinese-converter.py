import opencc
from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR

# Initialize OpenCC converter
converter = opencc.OpenCC('t2s.json')
traditional2simplified = opencc.OpenCC('t2s.json')
simplified2traditional = opencc.OpenCC('s2t.json')

t2s = True  # Set to True to convert from Traditional Chinese to Simplified Chinese, or False to convert from Simplified Chinese to Traditional Chinese

# Load PPT file
ppt = Presentation('input.pptx')  # Replace with the path to your input PPT file

# Iterate through each slide in the PPT
for slide in ppt.slides:
    # Iterate through each shape in the slide
    for shape in slide.shapes:
        # Check if the shape contains text
        if shape.has_text_frame and shape.text_frame.text.strip():
            # Convert the text from traditional to simplified Chinese
            text_frame = shape.text_frame
            if t2s:
                text = traditional2simplified.convert(text_frame.text)
            else:
                text = simplified2traditional.convert(text_frame.text)

            # Create a new text frame with the converted text
            new_frame = slide.shapes.add_textbox(shape.left, shape.top, shape.width, shape.height).text_frame
            new_frame.text = text

            # Set the text formatting properties of the new text frame to match the original text frame
            for p in text_frame.paragraphs:
                new_p = new_frame.add_paragraph()
                new_p.text = p.text
                new_p.font.name = p.font.name
                new_p.font.size = p.font.size
                new_p.font.bold = p.font.bold
                new_p.font.italic = p.font.italic
                new_p.line_spacing = p.line_spacing
                new_p.space_before = p.space_before
                new_p.space_after = p.space_after
                new_p.level = p.level
                new_p.font.language_id = p.font.language_id
                
                # Set the font color of the new text frame only if the font color of the original text frame is set
                if p.font.color.type != None:
                    new_p.font.color.rgb = p.font.color.rgb

            # Remove the original text frame
            slide.shapes._spTree.remove(shape._element)

# Save the modified PPT file
ppt.save('output.pptx')  # Replace with the path to your output PPT file
